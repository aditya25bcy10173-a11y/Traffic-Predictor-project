import sys
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
# Change slide dimensions to widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(10, 17, 40)
LIGHT_BLUE = RGBColor(0, 212, 255)
SAFFRON = RGBColor(255, 153, 51)
GREEN = RGBColor(19, 136, 8)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(240, 242, 245)
DARK_GRAY = RGBColor(60, 60, 60)

# Helper function to add a background color to a slide
def set_slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# Helper function to add slide title
def add_slide_title(slide, text, color=DARK_BLUE):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Arial"
    return txBox

# Create Title Slide (Slide 1)
slide_layout = prs.slide_layouts[6] # blank layout
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, DARK_BLUE)

# Add a glowing top line with Saffron/White/Green (Indian Flag touch)
top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
top_line.fill.solid()
top_line.fill.fore_color.rgb = SAFFRON
top_line.line.color.rgb = SAFFRON

# Title text box
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "BENGALURU INTELLIGENT MOBILITY\nCOMMAND CENTER"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.LEFT
p.font.name = "Arial"

# Subtitle
p2 = tf.add_paragraph()
p2.text = "Predictive Event-Driven Traffic Congestion & Tactical Resource Dispatch Platform"
p2.font.size = Pt(20)
p2.font.color.rgb = LIGHT_BLUE
p2.space_before = Pt(20)
p2.font.name = "Arial"

# Footer details (Indian Built / NIC / Govt of Karnataka)
footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.333), Inches(1.5))
tf_foot = footer_box.text_frame
p_foot = tf_foot.paragraphs[0]
p_foot.text = "Joint Smart Cities Mission Initiative: Directorate of Urban Land Transport (DULT) & Bengaluru Traffic Police (BTP)"
p_foot.font.size = Pt(14)
p_foot.font.bold = True
p_foot.font.color.rgb = SAFFRON
p_foot.font.name = "Arial"

p_foot2 = tf_foot.add_paragraph()
p_foot2.text = "Designed, Developed & Managed by National Informatics Centre (NIC) • Built in India 🇮🇳"
p_foot2.font.size = Pt(13)
p_foot2.font.color.rgb = RGBColor(200, 200, 200)
p_foot2.space_before = Pt(8)
p_foot2.font.name = "Arial"


# Slide 2: Project Overview & Objectives
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, LIGHT_GRAY)
add_slide_title(slide, "1. Project Overview & Objectives")

# Left Column: The Problem
left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.0))
tf_left = left_box.text_frame
tf_left.word_wrap = True
p = tf_left.paragraphs[0]
p.text = "The Challenge in Bengaluru"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

points = [
    "High density of planned (IPL matches, rallies) and unplanned (waterlogging, breakdowns) events causes localized, critical traffic bottlenecks.",
    "Traditional traffic control is reactive and lacks quantitative predictions of congestion duration.",
    "Manual marshal and police deployments are often misallocated due to lack of real-time predictive insights."
]
for pt in points:
    p_pt = tf_left.add_paragraph()
    p_pt.text = "• " + pt
    p_pt.font.size = Pt(16)
    p_pt.font.color.rgb = DARK_GRAY
    p_pt.space_before = Pt(14)
    p_pt.font.name = "Arial"

# Right Column: The Solution
right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.0))
tf_right = right_box.text_frame
tf_right.word_wrap = True
p = tf_right.paragraphs[0]
p.text = "The Intelligent Mobility Solution"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = GREEN

sol_points = [
    "Machine Learning Forecasting: Employs a state-of-the-art hybrid ensemble model (PyTorch Tabular ResNet + CatBoost) to predict duration of event impact.",
    "Tactical Dispatch recommendations: Recommends precise number of officers, marshals, towing rigs, barricade sizes, and diversion paths.",
    "NIC Command Portal: A premium, bilingual (English/Kannada) secure web portal featuring real-time GIS mapping, Adhira AI Assistant, and feedback loops."
]
for pt in sol_points:
    p_pt = tf_right.add_paragraph()
    p_pt.text = "• " + pt
    p_pt.font.size = Pt(16)
    p_pt.font.color.rgb = DARK_GRAY
    p_pt.space_before = Pt(14)
    p_pt.font.name = "Arial"


# Slide 3: System Architecture
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, LIGHT_GRAY)
add_slide_title(slide, "2. System Architecture & Flow")

# Left Column: Steps
flow_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.933), Inches(5.0))
tf_flow = flow_box.text_frame
tf_flow.word_wrap = True

steps = [
    ("Step 1: Input Event Parameters", "Incident type, cause, coordinates (lat/lon), priority level, road closure requirements, responsible police station, and temporal parameters (time/date)."),
    ("Step 2: AI Ensemble Inference", "Combines predictions from a PyTorch Deep Neural Network and CatBoost Regressor to estimate the resolution and congestion duration."),
    ("Step 3: Tactical Dispatch Recommendations", "Generates barricading, towing, diversion, and marshal requirement orders based on the predicted impact duration."),
    ("Step 4: NIC Command Portal & Dispatch", "Displays orders in the bilingual command center, prints official signed orders, and logs actual feedback to the database.")
]

for title, desc in steps:
    p_t = tf_flow.add_paragraph()
    p_t.text = "✔ " + title
    p_t.font.bold = True
    p_t.font.size = Pt(18)
    p_t.font.color.rgb = DARK_BLUE
    p_t.space_before = Pt(10)
    
    p_d = tf_flow.add_paragraph()
    p_d.text = "    " + desc
    p_d.font.size = Pt(15)
    p_d.font.color.rgb = DARK_GRAY
    p_d.space_before = Pt(4)


# Slide 4: Neural-Boosted Machine Learning Model
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, LIGHT_GRAY)
add_slide_title(slide, "3. Neural-Boosted ML Forecasting Model")

# Text details
desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(6.0), Inches(5.0))
tf_desc = desc_box.text_frame
tf_desc.word_wrap = True

p = tf_desc.paragraphs[0]
p.text = "Predicting Log-Duration: log1p(duration_minutes)"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

ml_points = [
    "PyTorch Tabular ResNet (Weight = 0.10): Captures complex, high-dimensional non-linear interactions across spatial and temporal parameters.",
    "CatBoost Regressor (Weight = 0.90): Handles categorical attributes (corridors, stations, causes) with optimal gradient boosting speeds.",
    "Ensemble Blend (0.10*PyTorch + 0.90*CatBoost): Yields validation R² of 0.3533, outperforming either model individually.",
    "Feature Engineering: Utilizes cyclic encoding (sine/cosine transformations) for spatial latitudes, longitudes, hours, and days of the week."
]
for pt in ml_points:
    p_pt = tf_desc.add_paragraph()
    p_pt.text = "• " + pt
    p_pt.font.size = Pt(15)
    p_pt.font.color.rgb = DARK_GRAY
    p_pt.space_before = Pt(10)

# Add a table on the right showing model metrics comparison
x, y, cx, cy = Inches(7.2), Inches(1.8), Inches(5.4), Inches(3.0)
shape = slide.shapes.add_table(4, 3, x, y, cx, cy)
table = shape.table

# Set Column Widths
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(1.6)
table.columns[2].width = Inches(1.6)

# Set Header Text
headers = ["Model Type", "Validation Loss (MSE)", "Validation R²"]
for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

# Set Data Rows
data = [
    ["PyTorch Tabular ResNet", "0.6412", "0.2990"],
    ["CatBoost Regressor", "0.5921", "0.3527"],
    ["Hybrid Ensemble Blend", "0.5916", "0.3533"]
]
for row_idx, row in enumerate(data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else RGBColor(245, 247, 250)
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(13)
        p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT


# Slide 5: The Recommender Engine
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, LIGHT_GRAY)
add_slide_title(slide, "4. Recommender Engine & Resource Rules")

# Left Column: Recommendation Rules Table
x, y, cx, cy = Inches(0.7), Inches(1.8), Inches(6.0), Inches(4.5)
shape = slide.shapes.add_table(5, 3, x, y, cx, cy)
table = shape.table
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(1.9)
table.columns[2].width = Inches(1.9)

rec_headers = ["Congestion Duration", "Police Officers", "Civil Marshals"]
for col_idx, text in enumerate(rec_headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    p = cell.text_frame.paragraphs[0]
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.size = Pt(13)
    p.alignment = PP_ALIGN.CENTER

rec_data = [
    ["Short (<30 mins)", "1 - 2 Officers", "0 - 1 Marshals"],
    ["Medium (30 - 60 mins)", "2 - 4 Officers", "2 - 3 Marshals"],
    ["Long (60 - 180 mins)", "4 - 8 Officers", "4 - 8 Marshals"],
    ["Critical (>180 mins)", "8 - 12 Officers", "8 - 15 Marshals"]
]
for row_idx, row in enumerate(rec_data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 0 else RGBColor(245, 247, 250)
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER

# Right Column: Tactical Strategy Details
strategy_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.633), Inches(5.0))
tf_strat = strategy_box.text_frame
tf_strat.word_wrap = True

p = tf_strat.paragraphs[0]
p.text = "Deployment & Diversion Tactics"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GREEN

strat_points = [
    "Barricading Strategies: Scaled based on predicted congestion severity. Suggests 'Light cones' for minor blocks, leading to 'Heavy full barricades' and physical lane blockages for critical events.",
    "Diversion Routes: Automatically extracts neighboring sectors and junctions to formulate localized diversion rules (e.g. diverting Silk Board rainfall congestion via HSR Layout sector loops).",
    "Towing Dispatch: Suggests 1 to 2 heavy-duty towing cranes standby for unplanned accidents and commercial truck breakdowns in major corridors (Outer Ring Road, Hebbal)."
]
for pt in strat_points:
    p_pt = tf_strat.add_paragraph()
    p_pt.text = "✔ " + pt
    p_pt.font.size = Pt(14)
    p_pt.font.color.rgb = DARK_GRAY
    p_pt.space_before = Pt(12)


# Slide 6: Bilingual NIC Command Portal (UI/UX)
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, LIGHT_GRAY)
add_slide_title(slide, "5. Bilingual NIC Command Portal (UI/UX)")

portal_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.933), Inches(5.0))
tf_portal = portal_box.text_frame
tf_portal.word_wrap = True

p = tf_portal.paragraphs[0]
p.text = "Official Government Visual Style Overhaul"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

ui_points = [
    "Bilingual Mode Switching: Dynamic toggle switches the entire application layout, text labels, dropdowns, and Adhira AI responses between English and Kannada (ಕನ್ನಡ) on the fly, with English set as default.",
    "Accessibility & Indian Branding: Features standard GOV.IN accessibility options (Screen Reader, font size controls A+/A/A-) alongside a thin horizontal Saffron/White/Green tricolor band and inline national flag icon.",
    "Interactive Chatbox Interface: Redesigns the Adhira AI Assistant into an overlay dialogue box with speech bubbles, Nic badges, avatars, and 3 quick-command chips (IPL Match, Silk Board Flood, Peenya Anomaly).",
    "Command Dispatch Order: Outputs forecasts in a formal, authenticated government memorandum layout containing simulated barcodes, approval stamps, IPS signatures, and dispatch registers.",
    "Simple Map GIS Integrations: Clean 2D GIS overlay maps plotted in real-time using historical coordinates from flip.csv."
]
for pt in ui_points:
    p_pt = tf_portal.add_paragraph()
    p_pt.text = "• " + pt
    p_pt.font.size = Pt(15)
    p_pt.font.color.rgb = DARK_GRAY
    p_pt.space_before = Pt(10)


# Slide 7: Closed-Loop Operator Feedback Database
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, LIGHT_GRAY)
add_slide_title(slide, "6. Closed-Loop Operator Feedback Loop")

loop_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.933), Inches(5.0))
tf_loop = loop_box.text_frame
tf_loop.word_wrap = True

p = tf_loop.paragraphs[0]
p.text = "How the Feedback Engine Works"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

db_points = [
    "Step 1: Logging Recommendations: The moment the operator runs a forecast in the 'AI Event Planner' tab, the system generates a unique Event ID (e.g. EVT0619093005) and logs predicted metrics (duration, marshals, officers) in the SQLite feedback.db database.",
    "Step 2: Dispatch Confirmation: The operator dispatch command notifies divisions, sending official reports to field units.",
    "Step 3: Operator Override & Review: In the 'Control Center' tab, the operator select boxes any logged Event ID to log post-incident observations. The form prompts for actual duration, actual marshals, officers deployed, and operator remarks.",
    "Step 4: Dynamic SQLite Logging: Actual observations are written back into feedback.db next to initial forecast predictions.",
    "Step 5: Model Continuous Learning: Database entries serve as labeled training samples for future model retraining runs, closing the loop between machine learning forecasts and real-world tactical feedback."
]
for pt in db_points:
    p_pt = tf_loop.add_paragraph()
    p_pt.text = "✔ " + pt
    p_pt.font.size = Pt(15)
    p_pt.font.color.rgb = DARK_GRAY
    p_pt.space_before = Pt(10)


# Slide 8: Future Roadmap & Next Steps
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide, LIGHT_GRAY)
add_slide_title(slide, "7. Project Status & Future Roadmap")

road_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.933), Inches(5.0))
tf_road = road_box.text_frame
tf_road.word_wrap = True

p = tf_road.paragraphs[0]
p.text = "Deployment Status & Vision"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

road_points = [
    "Portal Deployments: The Streamlit NIC Command Portal is fully running on http://localhost:8501, verified with zero runtime compilation errors.",
    "Full-Stack API Integration: Replicated architecture in a decoupled React Vite + FastAPI full-stack application at D:\bengaluru-command-center\ ready for cloud scaling.",
    "Future Expansion Phase 1: Real-time CCTV ingestion using computer vision to automatically detect and flag incident causes (accidents, vehicle breakdowns) without operator manual inputs.",
    "Future Expansion Phase 2: Live GPS integration with Namma Metro signals and BMTC buses to dynamically increase bus shuttle frequencies when major corridors experience forecast delays."
]
for pt in road_points:
    p_pt = tf_road.add_paragraph()
    p_pt.text = "• " + pt
    p_pt.font.size = Pt(15)
    p_pt.font.color.rgb = DARK_GRAY
    p_pt.space_before = Pt(12)

# Save presentation
prs.save("bengaluru_command_center_presentation.pptx")
print("Presentation saved successfully!")
